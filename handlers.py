import os
import re
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def create_presentation(topic, content, lang="ru"):
    """
    Создаёт .pptx файл из markdown-контента.
    content — текст от ИИ с заголовками ### и пунктами - 
    Возвращает путь к готовому файлу.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Титульный слайд
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Фон титульного
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RgbColor(0x1A, 0x5F, 0x7A)  # тёмно-бирюзовый
    background.line.fill.background()

    # Заголовок
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = topic
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # Подзаголовок
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(12.3), Inches(1))
    tf2 = sub_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "BilimBot" if lang == "ru" else "BilimBot"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RgbColor(0xCC, 0xCC, 0xCC)
    p2.alignment = PP_ALIGN.CENTER

    # Парсим слайды из markdown
    slides_data = parse_markdown_slides(content)

    for slide_info in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Фон
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RgbColor(0xF5, 0xF5, 0xF5)  # светло-серый
        bg.line.fill.background()

        # Заголовок слайда
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.3), Inches(1))
        tf_title = title_shape.text_frame
        p_title = tf_title.paragraphs[0]
        p_title.text = slide_info["title"]
        p_title.font.size = Pt(32)
        p_title.font.bold = True
        p_title.font.color.rgb = RgbColor(0x1A, 0x5F, 0x7A)
        p_title.alignment = PP_ALIGN.LEFT

        # Пункты
        if slide_info["points"]:
            body_shape = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12), Inches(5.5))
            tf_body = body_shape.text_frame
            tf_body.word_wrap = True

            for i, point in enumerate(slide_info["points"]):
                if i == 0:
                    p = tf_body.paragraphs[0]
                else:
                    p = tf_body.add_paragraph()
                p.text = f"• {point}"
                p.font.size = Pt(18)
                p.font.color.rgb = RgbColor(0x33, 0x33, 0x33)
                p.space_after = Pt(12)
                p.level = 0

    # Сохраняем
    tmp_dir = tempfile.gettempdir()
    filename = re.sub(r'[\\/*?:"<>|]', "_", topic)[:50] + ".pptx"
    filepath = os.path.join(tmp_dir, filename)
    prs.save(filepath)
    return filepath


def parse_markdown_slides(content):
    """
    Парсит markdown-ответ от ИИ.
    Ищет ### Заголовок и пункты - ...
    """
    slides = []
    current_slide = None

    for line in content.split("\n"):
        line = line.strip()
        if not line:
            continue

        # Заголовок слайда: ### Название
        if line.startswith("###"):
            if current_slide:
                slides.append(current_slide)
            current_slide = {"title": line.replace("###", "").strip(), "points": []}

        # Пункт: - текст
        elif line.startswith("-") or line.startswith("*"):
            if current_slide is not None:
                point = line[1:].strip()
                if point:
                    current_slide["points"].append(point)

        # Обычный текст — тоже добавляем как пункт, если уже есть слайд
        elif current_slide is not None and len(current_slide["points"]) < 8:
            current_slide["points"].append(line)

    if current_slide:
        slides.append(current_slide)

    # Если парсинг не сработал — создаём один слайд с полным текстом
    if not slides and content.strip():
        slides.append({"title": "Основное", "points": [content.strip()[:500]]})

    return slides
