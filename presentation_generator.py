from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import requests
import os
import re
from io import BytesIO
# ============================================================
# НАСТРОЙКИ
# ============================================================
WIKIMEDIA_API = "https://commons.wikimedia.org/w/api.php"
HEADERS = {
    "User-Agent": "BilimBot/1.0 educational Telegram bot"
}
# ============================================================
# БЕЗОПАСНОЕ ИМЯ ФАЙЛА
# ============================================================
def safe_filename(text):
    text = re.sub(
        r'[\\/*?:"<>|]',
        "",
        text
    )
    text = text.replace(" ", "_")
    return text[:35]
# ============================================================
# ПОИСК ФОТО В WIKIMEDIA COMMONS
# ============================================================
def search_wikimedia_image(query, index):
    """
    Бесплатный поиск изображения через Wikimedia Commons.
    API-ключ не нужен.
    """
    try:
        params = {
            "action": "query",
            "generator": "search",
            "gsrsearch": query,
            "gsrnamespace": 6,
            "gsrlimit": 10,
            "prop": "imageinfo",
            "iiprop": "url|mime|size",
            "iiurlwidth": 1000,
            "format": "json",
            "origin": "*"
        }
        response = requests.get(
            WIKIMEDIA_API,
            params=params,
            headers=HEADERS,
            timeout=15
        )
        if response.status_code != 200:
            print(
                "Wikimedia API error:",
                response.status_code
            )
            return None
        data = response.json()
        pages = (
            data
            .get("query", {})
            .get("pages", {})
        )
        if not pages:
            return None
        # Пробуем несколько результатов
        for number, page in enumerate(
            pages.values()
        ):
            imageinfo = page.get(
                "imageinfo",
                []
            )
            if not imageinfo:
                continue
            info = imageinfo[0]
            image_url = info.get(
                "thumburl"
            ) or info.get(
                "url"
            )
            mime = info.get(
                "mime",
                ""
            )
            if not image_url:
                continue
            # Нужны только обычные изображения
            if not mime.startswith("image/"):
                continue
            try:
                image_response = requests.get(
                    image_url,
                    headers=HEADERS,
                    timeout=15
                )
                if image_response.status_code != 200:
                    continue
                image_data = image_response.content
                if len(image_data) < 5000:
                    continue
                # Проверяем, что это действительно изображение
                try:
                    from PIL import Image
                    image = Image.open(
                        BytesIO(image_data)
                    )
                    image.verify()
                except Exception:
                    continue
                path = (
                    f"/tmp/"
                    f"bilim_slide_{index}_{number}.jpg"
                )
                # Pillow нужен для преобразования
                # PNG/WebP и других форматов в JPEG
                try:
                    from PIL import Image
                    image = Image.open(
                        BytesIO(image_data)
                    ).convert("RGB")
                    image.save(
                        path,
                        "JPEG",
                        quality=90
                    )
                except Exception:
                    with open(
                        path,
                        "wb"
                    ) as file:
                        file.write(
                            image_data
                        )
                print(
                    "Image found:",
                    query
                )
                return path
            except Exception as error:
                print(
                    "Image download error:",
                    error
                )
                continue
    except Exception as error:
        print(
            "Wikimedia search error:",
            error
        )
    return None
# ============================================================
# СОЗДАНИЕ ПРЕЗЕНТАЦИИ
# ============================================================
def create_presentation(
    topic,
    content_text,
    lang="ru"
):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    # ========================================================
    # ПАРСИНГ
    # ========================================================
    lines = content_text.strip().split(
        "\n"
    )
    slides_data = []
    current = None
    for line in lines:
        line = line.strip()
        if not line:
            continue
        # Заголовок
        if line.startswith("#"):
            if current:
                slides_data.append(
                    current
                )
            title = re.sub(
                r"^#+\s*",
                "",
                line
            ).strip()
            current = {
                "title": title,
                "points": []
            }
        # Пункт
        elif (
            line.startswith("- ")
            or line.startswith("• ")
            or line.startswith("* ")
        ):
            if current:
                current["points"].append(
                    line[2:].strip()
                )
        # Текст без заголовка
        elif current is None:
            current = {
                "title": line,
                "points": []
            }
    if current:
        slides_data.append(
            current
        )
    # Если структура не распарсилась
    if not slides_data:
        slides_data = [
            {
                "title": topic,
                "points": [
                    content_text[:1000]
                ]
            }
        ]
    # ========================================================
    # ТИТУЛЬНЫЙ СЛАЙД
    # ========================================================
    slide = prs.slides.add_slide(
        blank_layout
    )
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        prs.slide_width,
        prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(
        25,
        55,
        140
    )
    bg.line.fill.background()
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    # Нижняя полоска
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(6.8),
        prs.slide_width,
        Inches(0.7)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(
        255,
        190,
        40
    )
    accent.line.fill.background()
    # Заголовок
    tb = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(2.2),
        Inches(12.3),
        Inches(1.8)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = topic
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(
        255,
        255,
        255
    )
    p.alignment = PP_ALIGN.CENTER
    # Подзаголовок
    sub = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(4.2),
        Inches(12.3),
        Inches(0.8)
    )
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = (
        "BilimBot — Школьный помощник"
    )
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(
        180,
        200,
        255
    )
    p.alignment = PP_ALIGN.CENTER
    # ========================================================
    # СЛАЙДЫ
    # ========================================================
    selected_slides = slides_data[:10]
    total = len(selected_slides)
    for idx, sdata in enumerate(
        selected_slides
    ):
        slide = prs.slides.add_slide(
            blank_layout
        )
        # ====================================================
        # ФОН
        # ====================================================
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            prs.slide_width,
            prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(
            248,
            249,
            252
        )
        bg.line.fill.background()
        spTree = slide.shapes._spTree
        sp = bg._element
        spTree.remove(sp)
        spTree.insert(2, sp)
        # ====================================================
        # ШАПКА
        # ====================================================
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            prs.slide_width,
            Inches(1.0)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(
            35,
            75,
            170
        )
        header.line.fill.background()
        ht = header.text_frame
        ht.word_wrap = True
        hp = ht.paragraphs[0]
        hp.text = sdata["title"]
        hp.font.size = Pt(24)
        hp.font.bold = True
        hp.font.color.rgb = RGBColor(
            255,
            255,
            255
        )
        # ====================================================
        # ЖЁЛТАЯ ЛИНИЯ
        # ====================================================
        yellow = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(1.0),
            prs.slide_width,
            Inches(0.06)
        )
        yellow.fill.solid()
        yellow.fill.fore_color.rgb = RGBColor(
            255,
            190,
            40
        )
        yellow.line.fill.background()
        # ====================================================
        # ТЕКСТ СЛЕВА
        # ====================================================
        content = slide.shapes.add_textbox(
            Inches(0.4),
            Inches(1.25),
            Inches(8.35),
            Inches(5.9)
        )
        tf = content.text_frame
        tf.word_wrap = True
        for point in sdata[
            "points"
        ][:8]:
            p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(
                35,
                35,
                35
            )
            p.space_after = Pt(11)
            p.level = 0
        # ====================================================
        # ПОИСК ИЗОБРАЖЕНИЯ
        # ====================================================
        image_path = None
        # Сначала ищем по заголовку
        search_query = sdata["title"]
        # Если заголовок короткий
        if len(search_query) < 5:
            search_query = (
                f"{topic} "
                f"{search_query}"
            )
        # Убираем некоторые символы,
        # которые мешают поиску
        search_query = re.sub(
            r"[#*_]",
            "",
            search_query
        ).strip()
        print(
            f"[BilimBot] "
            f"Searching image for: "
            f"{search_query}"
        )
        image_path = search_wikimedia_image(
            search_query,
            idx
        )
        # Если не найдено —
        # пробуем общий запрос по теме
        if not image_path:
            print(
                "[BilimBot] "
                "Trying topic search..."
            )
            image_path = search_wikimedia_image(
                topic,
                idx
            )
        # ====================================================
        # ФОТО НАЙДЕНО
        # ====================================================
        if image_path:
            try:
                # Рамка
                frame = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(9.05),
                    Inches(1.3),
                    Inches(3.85),
                    Inches(4.8)
                )
                frame.fill.solid()
                frame.fill.fore_color.rgb = RGBColor(
                    230,
                    235,
                    245
                )
                frame.line.color.rgb = RGBColor(
                    200,
                    210,
                    230
                )
                # Фото
                slide.shapes.add_picture(
                    image_path,
                    Inches(9.2),
                    Inches(1.45),
                    width=Inches(3.55),
                    height=Inches(4.5)
                )
                print(
                    f"[BilimBot] "
                    f"Image added to slide "
                    f"{idx + 1}"
                )
            except Exception as error:
                print(
                    "[BilimBot] "
                    f"Could not add image: "
                    f"{error}"
                )
                image_path = None
        # ====================================================
        # ЕСЛИ ФОТО НЕ НАШЛОСЬ
        # ====================================================
        if not image_path:
            decor = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(9.5),
                Inches(2.0),
                Inches(3.0),
                Inches(3.0)
            )
            decor.fill.solid()
            decor.fill.fore_color.rgb = RGBColor(
                230,
                235,
                245
            )
            decor.line.color.rgb = RGBColor(
                200,
                210,
                230
            )
            dt = decor.text_frame
            dt.word_wrap = True
            dp = dt.paragraphs[0]
            dp.text = "📷"
            dp.font.size = Pt(48)
            dp.alignment = PP_ALIGN.CENTER
        # ====================================================
        # НОМЕР
        # ====================================================
        num = slide.shapes.add_textbox(
            Inches(11.8),
            Inches(7.05),
            Inches(1.2),
            Inches(0.35)
        )
        nt = num.text_frame
        np = nt.paragraphs[0]
        np.text = (
            f"{idx + 1} / {total}"
        )
        np.font.size = Pt(11)
        np.font.color.rgb = RGBColor(
            150,
            150,
            150
        )
        np.alignment = PP_ALIGN.RIGHT
        # ====================================================
        # УДАЛЕНИЕ ВРЕМЕННОГО ФАЙЛА
        # ====================================================
        try:
            if (
                image_path
                and os.path.exists(
                    image_path
                )
            ):
                os.remove(
                    image_path
                )
        except Exception:
            pass
    # ========================================================
    # ФИНАЛЬНЫЙ СЛАЙД
    # ========================================================
    slide = prs.slides.add_slide(
        blank_layout
    )
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        prs.slide_width,
        prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(
        25,
        55,
        140
    )
    bg.line.fill.background()
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    end = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(2.8),
        Inches(12.3),
        Inches(1.5)
    )
    tf = end.text_frame
    p = tf.paragraphs[0]
    p.text = "Спасибо за внимание!"
    p.font.size = Pt(46)
    p.font.bold = True
    p.font.color.rgb = RGBColor(
        255,
        255,
        255
    )
    p.alignment = PP_ALIGN.CENTER
    sub = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(4.3),
        Inches(12.3),
        Inches(0.8)
    )
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = "Сгенерировано BilimBot"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(
        180,
        200,
        255
    )
    p.alignment = PP_ALIGN.CENTER
    # ========================================================
    # СОХРАНЕНИЕ
    # ========================================================
    safe = safe_filename(topic)
    output = (
        f"/tmp/pres_{safe}.pptx"
    )
    prs.save(output)
    print(
        f"[BilimBot] "
        f"Presentation saved: {output}"
    )
    return output
